import logging
from dataclasses import dataclass

from pptx import Presentation

# from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Pt

from config import (
    DICT_LIST_CONTENT_FONT_SIZE_PT,
    EMU_PER_INCH,
    EMU_PER_PT,
    INDUSTRY_CONTINUATION_UPLIFT_INCH,
    MULTILINE_CONTENT_FONT_SIZE_PT,
    SIMPLE_CONTENT_FONT_SIZE_PT,
    SIMPLE_LIST_CONTENT_FONT_SIZE_PT,
    TABLE_HEADER_BG_COLOR,
    TABLE_HEADER_HEIGHT_PT,
    TABLE_HEADER_TEXT_COLOR,
    TABLE_LINE_HEIGHT_PT,
    TABLE_PARAGRAPH_FONT_SIZE_PT,
)
from helpers.exceptions import TemplateError
from helpers.utils import _remove_shape_and_get_bbox, estimate_row_height, unwrap_first_data


@dataclass
class TableDimensions:
    """Holds calculated dimensions for table layout."""

    total_height_pt: float
    header_height_pt: float
    content_height_pt: float
    line_height_pt: float
    width_pt: float
    column_width_pt: float


def fill_industry_slides(prs: Presentation, payload: dict):
    """
    Replace the {{IndustryResearch}} token by locating the placeholder on any slide.
    Uses payload['title'], payload['headers'], payload['rows'] (after unwrapping payload['data'][0]).
    If headers/rows are missing, it logs a warning and leaves the presentation unchanged.
    """
    slide, placeholder_shape = _find_placeholder(prs, "{{IndustryResearch}}")
    if not placeholder_shape:
        raise TemplateError("Token '{{IndustryResearch}}' not found in any slide")

    # Normaliza payload (acepta wrapper con data[])
    try:
        payload_norm = unwrap_first_data(payload, "IndustryResearch")
    except Exception as exc:
        raise TemplateError(f"Invalid IndustryResearch payload: {exc}") from exc

    title_text = _set_slide_title(slide, payload_norm)

    headers, rows = _validate_and_extract_data(payload_norm)
    if not headers or not rows:
        logging.warning("IndustryResearch payload missing headers/rows; skipping slide rendering.")
        return

    # Remueve placeholder solo cuando hay datos válidos para no dejar la slide vacía.
    left, top, width, height = _remove_shape_and_get_bbox(placeholder_shape)
    layout = slide.slide_layout

    # Manually lift continuation slides to reclaim the missing title area.
    UPLIFT_EMU = int(INDUSTRY_CONTINUATION_UPLIFT_INCH * EMU_PER_INCH)
    cont_top = max(0, top - UPLIFT_EMU)
    cont_height = height + (top - cont_top)

    dimensions = _calculate_table_dimensions(width, height, len(headers))
    row_heights = _calculate_row_heights(rows, headers, dimensions)

    chunks = _partition_rows_into_chunks(rows, row_heights, dimensions.content_height_pt)

    for idx, chunk in enumerate(chunks):
        target_slide = _get_or_create_slide(prs, slide, layout, idx, title_text)
        # First chunk uses original bbox; continuations reclaim title space.
        top_use = top if idx == 0 else cont_top
        height_use = height if idx == 0 else cont_height
        table = _create_table(target_slide, chunk, headers, left, top_use, width, height_use)
        _format_table_header(table, headers, dimensions, width)
        _populate_table_data(table, chunk, headers)


# particionar filas en trozos que quepan en content_pt
def _partition_rows_into_chunks(rows: list, row_heights: list, available_height_pt: float) -> list | None:
    """
    Partitions rows into chunks that fit within the available height.

    Args:
        rows: List of row data
        row_heights: List of estimated heights for each row
        available_height_pt: Available height in points

    Returns:
        List of row chunks that fit within the height constraint
    """
    chunks = []
    current_index = 0
    total_rows = len(rows)

    while current_index < total_rows:
        chunk_height_used = 0.0
        chunk_end_index = current_index

        # Add rows to chunk while they fit within available height
        while chunk_end_index < total_rows and chunk_height_used + row_heights[chunk_end_index] <= available_height_pt:
            chunk_height_used += row_heights[chunk_end_index]
            chunk_end_index += 1

        # Ensure at least one row is included in each chunk
        if chunk_end_index == current_index:
            chunk_end_index = current_index + 1

        chunks.append(rows[current_index:chunk_end_index])
        current_index = chunk_end_index

    return chunks


def _find_placeholder(prs: Presentation, token: str):
    """Find first slide and shape containing the token without removing it."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and token in shape.text_frame.text:
                return slide, shape
    return None, None


def _set_slide_title(slide, payload: dict) -> str:
    """Sets the slide title from payload data. Returns the text used."""
    if slide.shapes.title:
        slide.shapes.title.text = payload.get("title", "")
        return slide.shapes.title.text
    return ""


def _validate_and_extract_data(payload: dict) -> tuple[list, list]:
    """
    Validates and extracts headers and rows from payload.

    Returns:
        Tuple of (headers, rows) or ([], []) if validation fails
    """
    headers = payload.get("headers", [])
    rows = payload.get("rows", [])

    if not headers or not rows:
        return [], []

    return headers, rows


def _calculate_table_dimensions(width: int, height: int, column_count: int) -> TableDimensions:
    """
    Calculates table dimensions for layout planning.

    Args:
        width: Table width in EMU units
        height: Table height in EMU units
        column_count: Number of columns in the table

    Returns:
        TableDimensions object with calculated values
    """
    total_height_pt = height / EMU_PER_PT
    header_height_pt = TABLE_HEADER_HEIGHT_PT
    content_height_pt = total_height_pt - header_height_pt
    line_height_pt = TABLE_LINE_HEIGHT_PT

    width_pt = width / EMU_PER_PT
    column_width_pt = width_pt / column_count

    return TableDimensions(
        total_height_pt=total_height_pt,
        header_height_pt=header_height_pt,
        content_height_pt=content_height_pt,
        line_height_pt=line_height_pt,
        width_pt=width_pt,
        column_width_pt=column_width_pt,
    )


def _calculate_row_heights(rows: list, headers: list, dimensions: TableDimensions) -> list[float]:
    """
    Calculates estimated height for each row.

    Args:
        rows: List of row data
        headers: List of column headers
        dimensions: Table dimensions object

    Returns:
        List of estimated heights for each row
    """
    return [estimate_row_height(row, headers, dimensions.line_height_pt, dimensions.column_width_pt) for row in rows]


def _get_or_create_slide(prs: Presentation, original_slide, layout, chunk_index: int, title_text: str):
    """
    Gets the original slide for the first chunk or creates a new slide for subsequent chunks.

    Args:
        prs: PowerPoint presentation object
        original_slide: The original slide to use for the first chunk
        layout: Slide layout to use for new slides
        chunk_index: Index of the current chunk (0-based)

    Returns:
        Slide object to use for the current chunk
    """
    if chunk_index == 0:
        return original_slide

    new_slide = prs.slides.add_slide(layout)
    _clear_placeholders(new_slide)
    # No title on continuation pages to avoid vertical gap; table stays flush to bbox.
    return new_slide


def _clear_placeholders(slide):
    """Remove all placeholders from a slide (title/content) to free vertical space."""
    for shape in list(slide.shapes):
        if getattr(shape, "is_placeholder", False):
            slide.shapes._spTree.remove(shape._element)


def _create_table(slide, chunk: list, headers: list, left: int, top: int, width: int, height: int):
    """
    Creates a table on the specified slide with the given dimensions.

    Args:
        slide: PowerPoint slide object
        chunk: List of row data for this table
        headers: List of column headers
        left, top, width, height: Table position and dimensions in EMU units

    Returns:
        Table object
    """
    rows_count = len(chunk) + 1  # +1 for header row
    cols_count = len(headers)

    return slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table


def _format_table_header(table, headers: list, dimensions: TableDimensions, width: int) -> None:
    """
    Formats the header row of the table with styling and content.

    Args:
        table: PowerPoint table object
        headers: List of column headers
        dimensions: Table dimensions object
        width: Total table width in EMU units
    """
    # Set header row height
    table.rows[0].height = Pt(dimensions.header_height_pt)

    # Set column widths
    cols_count = len(headers)
    for col in table.columns:
        col.width = width // cols_count

    # Format header cells
    for col_index, header_text in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.text = header_text

        # Cell background
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG_COLOR

        # Text formatting
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER

        for run in paragraph.runs:
            run.font.size = Pt(TABLE_PARAGRAPH_FONT_SIZE_PT)
            run.font.bold = True
            run.font.color.rgb = TABLE_HEADER_TEXT_COLOR

        cell.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def _populate_table_data(table, chunk: list, headers: list) -> None:
    """
    Populates the table with data from the chunk.

    Args:
        table: PowerPoint table object
        chunk: List of row data
        headers: List of column headers
    """
    for row_index, entry in enumerate(chunk, start=1):
        for col_index, header in enumerate(headers):
            cell = table.cell(row_index, col_index)
            text_frame = cell.text_frame
            text_frame.clear()
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

            value = entry.get(header, "")
            _format_cell_content(text_frame, value)


def _format_cell_content(text_frame, value) -> None:
    """
    Formats the content of a table cell based on the value type.

    Args:
        text_frame: PowerPoint text frame object
        value: The value to format and add to the cell
    """
    if isinstance(value, list) and value and isinstance(value[0], dict):
        _add_dict_list_content(text_frame, value)
    elif isinstance(value, list):
        _add_simple_list_content(text_frame, value)
    elif isinstance(value, str) and "\n" in value:
        _add_multiline_content(text_frame, value)
    else:
        _add_simple_content(text_frame, str(value))


def _add_dict_list_content(text_frame, dict_list: list) -> None:
    """Adds content for a list of dictionaries."""
    for item in dict_list:
        line = "; ".join(f"{k}: {v}" for k, v in item.items())
        paragraph = text_frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.size = Pt(DICT_LIST_CONTENT_FONT_SIZE_PT)


def _add_simple_list_content(text_frame, item_list: list) -> None:
    """Adds content for a simple list with bullet points."""
    for item in item_list:
        paragraph = text_frame.add_paragraph()
        paragraph.text = f"• {item}"
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.size = Pt(SIMPLE_LIST_CONTENT_FONT_SIZE_PT)


def _add_multiline_content(text_frame, text: str) -> None:
    """Adds content for multiline text."""
    for line in text.splitlines():
        paragraph = text_frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.size = Pt(MULTILINE_CONTENT_FONT_SIZE_PT)


def _add_simple_content(text_frame, text: str) -> None:
    """Adds simple text content."""
    paragraph = text_frame.add_paragraph()
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.LEFT
    for run in paragraph.runs:
        run.font.size = Pt(SIMPLE_CONTENT_FONT_SIZE_PT)
