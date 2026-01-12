import json
import math
import os
import re
from datetime import datetime

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt


def _set_font_size(run, size_pt=8, bold=False, color=None):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def _find_shape_with_token(prs: Presentation, token: str):
    """Return the first shape in the presentation that contains the token."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and token in shape.text_frame.text:
                return slide, shape
    return None, None


def _replace_token_in_shape_text(shape, token: str, value: str):
    """Replace the token with value while preserving the rest of the shape text."""
    tf = shape.text_frame
    # Replace across all paragraphs/runs
    for p in tf.paragraphs:
        new_text = p.text.replace(token, value)
        if new_text != p.text:
            p.text = new_text


def _replace_company_name_everywhere(prs: Presentation, name: str):
    """Replace {{CompanyName}} tokens across all shapes in the presentation."""
    tokens = ["{{CompanyName}}", "{{ CompanyName }}"]
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for t in tokens:
                    if t in shape.text_frame.text:
                        _replace_token_in_shape_text(shape, t, name)


def _remove_shape_and_get_bbox(shape):
    """Remove the shape and return its bounding box for reuse (tables, etc.)."""
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    shape._element.getparent().remove(shape._element)
    return left, top, width, height


def _add_section_header(tf, title: str, size: int = 18):
    p = tf.add_paragraph()
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    for run in p.runs:
        _set_font_size(run, size_pt=size, bold=True)


def _add_bullet(tf, text: str, level: int = 0, size: int = 14):
    p = tf.add_paragraph()
    p.text = f"• {text}" if not text.startswith("•") else text
    p.alignment = PP_ALIGN.LEFT
    p.level = level
    for run in p.runs:
        _set_font_size(run, size_pt=size)


# def _fmt_currency(n: int | float, prefix="$"):
#    try:
#        return f"{prefix}{int(n):,}".replace(",", ".")  # separador miles estilo LATAM opcional
#    except Exception:
#        return str(n)


def _norm(s: str) -> str:
    return re.sub(r"\s+", " ", str(s or "")).strip().lower()


def _is_url(s: str) -> bool:
    return isinstance(s, str) and s.startswith(("http://", "https://"))


def _extract_urls(obj) -> list[str]:
    urls = []
    if isinstance(obj, dict):
        for k, v in obj.items():
            if _is_url(v):
                urls.append(v)
            elif isinstance(v, (dict, list)):
                urls.extend(_extract_urls(v))
    elif isinstance(obj, list):
        for it in obj:
            urls.extend(_extract_urls(it))
    return urls


def _parse_date(s: str) -> str:
    """Return a readable date if ISO/YYYY-MM-DD/YYYY-MM/YYYY is recognized."""
    if not isinstance(s, str) or not s.strip():
        return ""
    t = s.strip()
    # Minimal normalizations
    try:
        # ISO completo
        return datetime.fromisoformat(t).strftime("%B %d, %Y")
    except Exception:
        pass
    # YYYY-MM
    m = re.match(r"^\d{4}-\d{2}$", t)
    if m:
        try:
            return datetime.fromisoformat(t + "-01").strftime("%B %d, %Y")
        except Exception:
            return t
    # YYYY
    if re.match(r"^\d{4}$", t):
        try:
            return datetime(int(t), 12, 31).strftime("%B %d, %Y")
        except Exception:
            return t
    return t


def _parse_number(x):
    """Try to coerce to a number (accepts '97.69 billion', '$97,690,000,000', etc.)."""
    if x is None:
        return None
    if isinstance(x, (int, float)):
        return float(x)
    s = str(x)
    # billion/million textuales
    m = re.match(r"^\s*\$?\s*([\d\.,]+)\s*(billion|million)\b", s, re.I)
    if m:
        base = (
            float(m.group(1).replace(".", "").replace(",", "."))
            if "." not in m.group(1)
            else float(m.group(1).replace(",", ""))
        )
        mult = 1e9 if m.group(2).lower() == "billion" else 1e6
        return base * mult
    # strip symbols
    s = re.sub(r"[^\d\.\-]", "", s)
    try:
        return float(s)
    except Exception:
        return None


def _fmt_billions_usd(n):
    if n is None:
        return ""
    try:
        return f"${n / 1e9:.2f} billion USD"
    except Exception:
        return str(n)


def _parse_percent(x):
    """Return float 0-100 if possible."""
    if x is None:
        return None
    if isinstance(x, (int, float)):
        return float(x)
    s = str(x).strip()
    s = s.replace(",", ".")
    m = re.match(r"^\s*([\-]?\d+(\.\d+)?)\s*%?\s*$", s)
    if not m:
        return None
    return float(m.group(1))


def _choose_link(*candidates):
    """Pick a preferred link (prioritize sec.gov when present)."""
    urls = []
    for c in candidates:
        if isinstance(c, str) and _is_url(c):
            urls.append(c)
        elif isinstance(c, (list, dict)):
            urls.extend([u for u in _extract_urls(c) if _is_url(u)])
    # prioridad sec.gov
    for u in urls:
        if "sec.gov" in u:
            return u
    return urls[0] if urls else None


def _find_in_dict(d: dict, key_synonyms: list[str]) -> dict | None:
    """Search current level for a key whose normalized name matches any synonym."""
    for k, v in d.items():
        nk = _norm(k)
        if any(s in nk for s in key_synonyms):
            return v if isinstance(v, dict) else {"value": v}
    return None


def _deep_find(d: dict, key_synonyms: list[str]) -> dict | None:
    """Recursive search when not found at the first level."""
    hit = _find_in_dict(d, key_synonyms)
    if hit is not None:
        return hit
    for _, v in d.items():
        if isinstance(v, dict):
            h = _deep_find(v, key_synonyms)
            if h is not None:
                return h
        elif isinstance(v, list):
            for it in v:
                if isinstance(it, dict):
                    h = _deep_find(it, key_synonyms)
                    if h is not None:
                        return h
    return None


def _get_first_str(d: dict, key_synonyms: list[str]) -> str:
    """Return the first string matching any synonym (depth 1)."""
    for k, v in d.items():
        if any(s in _norm(k) for s in key_synonyms):
            if isinstance(v, str):
                return v
    return ""


def _add_bullet_runs(tf, runs, level=0, size=14):
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.level = level
    for piece in runs:
        r = p.add_run()
        r.text = piece.get("text", "")
        _set_font_size(r, size_pt=size, bold=piece.get("bold", False))
        if piece.get("link"):
            r.hyperlink.address = piece["link"]
    return p


def estimate_row_height(entry: dict, keys: list, line_height_pt: int, col_width_pt: float) -> float:
    """Estimate height of the row (in points) based on text length and wrapping."""
    avg_char_width_pt = 6.0  # approx average character width at 10pt font
    chars_per_line = max(int(col_width_pt / avg_char_width_pt), 1)
    max_lines = 1
    for h in keys:
        val = entry.get(h, "")
        # count lines for this cell
        if isinstance(val, list) and val and isinstance(val[0], dict):
            # list of dicts: each dict becomes a line, wrapped
            lines_i = 0
            for item in val:
                s = "; ".join(f"{k}: {v}" for k, v in item.items())
                # wrap count
                lines_i += math.ceil(len(s) / chars_per_line)
        elif isinstance(val, list):
            # list of primitives
            lines_i = 0
            for item in val:
                s = str(item)
                lines_i += math.ceil(len(s) / chars_per_line)
        elif isinstance(val, str):
            # split existing newlines, then wrap
            lines_i = 0
            for line in val.splitlines():
                lines_i += math.ceil(len(line) / chars_per_line)
        else:
            # single primitive
            s = str(val)
            lines_i = math.ceil(len(s) / chars_per_line)
        max_lines = max(max_lines, lines_i)
    return max_lines * line_height_pt


def _load_json(data_or_path):
    if isinstance(data_or_path, dict):
        return data_or_path
    if isinstance(data_or_path, (str, os.PathLike)):
        with open(data_or_path, encoding="utf-8") as f:
            return json.load(f)
    raise TypeError("Expected dict or JSON file path")


def unwrap_first_data(payload: dict, label: str) -> dict:
    """
    Normalize payload structures that wrap content under payload['data'][0].

    Returns the inner dict if present, otherwise returns the payload itself.
    Raises ValueError when the structure is invalid or empty.
    """
    if not isinstance(payload, dict):
        raise ValueError(f"{label} must be an object")

    if "data" not in payload:
        return payload

    data = payload.get("data")
    if not isinstance(data, list):
        raise ValueError(f"{label}['data'] must be a list")
    if not data:
        raise ValueError(f"{label}['data'] is empty")

    first = data[0]
    if not isinstance(first, dict):
        raise ValueError(f"{label}['data'][0] must be an object")

    return first
