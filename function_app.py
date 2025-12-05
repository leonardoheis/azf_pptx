import json
import logging
import os
from datetime import datetime, timezone
from functools import lru_cache
from io import BytesIO

import azure.functions as func
from azure.core.exceptions import ResourceExistsError
from azure.storage.blob import BlobServiceClient
from pptx import Presentation

from company_research1 import fill_company_name_from_json, fill_company_research1
from company_research2 import fill_company_research2
from company_research3 import fill_company_research3
from config import AZ_STORAGE_CONN_STRING  # may be None in local env
from config import AZ_BLOB_CONTAINER_NAME, INPUT_TEMPLATE, get_next_output_filename
from helpers.exceptions import AppError, ValidationError
from industry_research import fill_industry_slides

app = func.FunctionApp(http_auth_level=func.AuthLevel.ANONYMOUS)


@lru_cache(maxsize=1)
def _init_clients():
    """Create and return (table_client, blob_service_client, container_client).

    Cached as a single tuple to avoid module-level mutable globals.
    """
    conn = _get_conn_string()

    # Table client
    # svc = TableServiceClient.from_connection_string(conn)
    # table_client = svc.get_table_client(table_name=AZ_BLOB_TABLE_NAME)
    # try:
    #     table_client.create_table()
    # except ResourceExistsError:
    #     logging.debug("Table already exists (create_table). Continuing.")
    # except Exception:
    #     logging.exception("Unexpected error creating table")
    #     raise

    # Blob client
    blob_service_client = BlobServiceClient.from_connection_string(conn)
    container_client = blob_service_client.get_container_client(container=AZ_BLOB_CONTAINER_NAME)
    try:
        container_client.create_container()
    except ResourceExistsError:
        logging.debug("Container already exists (create_container). Continuing.")
    except Exception:
        logging.exception("Unexpected error creating container")
        raise
    # , table_client,
    return blob_service_client, container_client


def _get_conn_string() -> str:
    """
    Prefer AZ_STORAGE_CONN_STRING; fallback to AzureWebJobsStorage.
    Validate it early but *don't* fail at import-time.
    """
    conn = AZ_STORAGE_CONN_STRING or os.getenv("AZ_STORAGE_CONN_STRING")
    if not conn:
        conn = os.getenv("AzureWebJobsStorage")  # may be 'UseDevelopmentStorage=true'
    if not conn or len(conn.strip()) == 0:
        raise RuntimeError("Storage connection string missing. Set AZ_STORAGE_CONN_STRING or AzureWebJobsStorage.")
    return conn


# def _ensure_resources():
#    """Compatibility shim retained for call sites; triggers client initialization."""
#    _init_clients()


# def _get_table_client():
#    table_client, _, _ = _init_clients()
#    return table_client


def _get_container_client():
    _, container_client = _init_clients()
    return container_client


def _validate_request_data(req_body: dict) -> tuple[dict, dict, dict, dict, str | None]:
    """
    Validate that all required fields are present and of correct type.

    Returns:
        Tuple of (company_data1, company_data2, company_data3, industry_data, error_message)
        error_message is None if validation passes.
    """
    required = ["CompanyResearchData1", "CompanyResearchData2", "CompanyResearchData3", "IndustryResearch"]
    missing = [k for k in required if k not in req_body]
    if missing:
        raise ValidationError(f"Missing required files: {', '.join(missing)}")

    company_data1 = req_body["CompanyResearchData1"]
    company_data2 = req_body["CompanyResearchData2"]
    company_data3 = req_body["CompanyResearchData3"]
    industry_data = req_body["IndustryResearch"]

    for name, obj in [
        ("CompanyResearchData1", company_data1),
        ("CompanyResearchData2", company_data2),
        ("CompanyResearchData3", company_data3),
        ("IndustryResearch", industry_data),
    ]:
        if not isinstance(obj, dict):
            raise ValidationError(f"Field '{name}' must be a valid JSON object")

    return company_data1, company_data2, company_data3, industry_data, None


def _log_received_data_keys(company_data1: dict, company_data2: dict, company_data3: dict, industry_data: dict) -> None:
    """Log the keys received in each data dictionary."""
    logging.info("Keys1=%s", list(company_data1.keys()))
    logging.info("Keys2=%s", list(company_data2.keys()))
    logging.info("Keys3=%s", list(company_data3.keys()))
    logging.info("KeysIndustry=%s", list(industry_data.keys()))


def _create_processed_data_summary(
    company_data1: dict, company_data2: dict, company_data3: dict, industry_data: dict
) -> dict:
    """Create a summary dictionary of processed data statistics."""
    return {
        "company_data1_processed": len(company_data1),
        "company_data2_processed": len(company_data2),
        "company_data3_processed": len(company_data3),
        "industry_data_processed": len(industry_data),
        "total_fields": len(company_data1) + len(company_data2) + len(company_data3) + len(industry_data),
        "timestamp": datetime.now(timezone.utc).isoformat(),
        "output_location": None,
    }


def _build_presentation(
    company_data1: dict, company_data2: dict, company_data3: dict, industry_data: dict
) -> Presentation:
    """
    Build a PowerPoint presentation from the provided data.

    Returns:
        Presentation object populated with the data.
    """
    prs = Presentation(INPUT_TEMPLATE)

    fill_company_name_from_json(prs, company_data1)
    fill_company_research1(prs, company_data1)
    fill_company_research2(prs, company_data2)
    fill_company_research3(prs, company_data3)
    fill_industry_slides(prs, prs.slides[len(prs.slides) - 1], industry_data)

    return prs


def _save_presentation_to_buffer(prs: Presentation) -> BytesIO:
    """Save a presentation to a BytesIO buffer."""
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _upload_presentation_to_blob(buf: BytesIO, output_filename: str) -> str:
    """
    Upload presentation buffer to Azure Blob Storage.

    Returns:
        Full URL of the uploaded blob.
    """
    container_client = _get_container_client()
    blob_client = container_client.get_blob_client(blob=output_filename)
    blob_client.upload_blob(buf.getvalue(), overwrite=True)

    account = blob_client.account_name
    blob_url = f"https://{account}.blob.core.windows.net/{AZ_BLOB_CONTAINER_NAME}/{output_filename}"
    logging.info("Presentation saved: %s", blob_url)

    return blob_url


def _create_table_entity(
    company_data1: dict, company_data2: dict, company_data3: dict, industry_data: dict, processed_data: dict
) -> dict:
    """Create a table entity for logging processed data."""
    return {
        "PartitionKey": "processed_files",
        "RowKey": datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S"),
        "CompanyData1Keys": json.dumps(list(company_data1.keys())),
        "CompanyData2Keys": json.dumps(list(company_data2.keys())),
        "CompanyData3Keys": json.dumps(list(company_data3.keys())),
        "IndustryDataKeys": json.dumps(list(industry_data.keys())),
        "ProcessedData": json.dumps(processed_data),
        "Timestamp": datetime.now(timezone.utc).isoformat(),
    }


# def _log_to_table(entity: dict) -> None:
#    """Log an entity to Azure Table Storage."""
#    table_client = _get_table_client()
#    table_client.create_entity(entity=entity)


def _create_success_response(
    processed_data: dict,
    company_data1: dict,
    company_data2: dict,
    company_data3: dict,
    industry_data: dict,
    output_filename: str,
) -> dict:
    """Create the success response payload."""
    return {
        "status": "success",
        "message": "Processed files and saved PPTX to Blob",
        "data": processed_data,
        "files_received": {
            "CompanyResearchData1_size": len(company_data1),
            "CompanyResearchData2_size": len(company_data2),
            "CompanyResearchData3_size": len(company_data3),
            "IndustryResearch_size": len(industry_data),
        },
        "output_file": {
            "filename": output_filename,
            "blob_path": output_filename,
            "container": AZ_BLOB_CONTAINER_NAME,
            "full_url": processed_data.get("output_location", "Not available"),
        },
    }


def _create_error_response(error_message: str, status_code: int) -> func.HttpResponse:
    """Create an error HTTP response."""
    return func.HttpResponse(
        json.dumps({"error": error_message, "status": "error"}),
        status_code=status_code,
        mimetype="application/json",
    )


@app.route(route="agent_httptrigger")
def agent_httptrigger(req: func.HttpRequest) -> func.HttpResponse:
    """
    Receives 4 JSON files and produces a PPTX, uploads to Blob, and logs to Tables.
    """
    try:
        req_body = req.get_json()

        # Validate request data
        # Validation errors are now raised as exceptions
        company_data1, company_data2, company_data3, industry_data, _ = _validate_request_data(req_body)

        # Log received data
        _log_received_data_keys(company_data1, company_data2, company_data3, industry_data)

        # Create processing summary
        processed_data = _create_processed_data_summary(company_data1, company_data2, company_data3, industry_data)

        # Build and upload presentation
        current_output_file = get_next_output_filename()
        try:
            prs = _build_presentation(company_data1, company_data2, company_data3, industry_data)
            buf = _save_presentation_to_buffer(prs)
            blob_url = _upload_presentation_to_blob(buf, current_output_file)
            processed_data["output_location"] = blob_url
        except Exception as e:
            logging.exception("Error building/uploading PPTX: %s", e)

        # Log to Table Storage
        # try:
        #     entity = _create_table_entity(company_data1, company_data2, company_data3, industry_data, processed_data)
        #     _log_to_table(entity)
        # except Exception as e:
        #     logging.exception("Error storing data in table: %s", e)

        # Create and return success response
        response_data = _create_success_response(
            processed_data, company_data1, company_data2, company_data3, industry_data, current_output_file
        )
        return func.HttpResponse(
            json.dumps(response_data),
            status_code=200,
            mimetype="application/json",
        )

    except ValueError as e:
        return _create_error_response(f"Invalid JSON: {e}", 400)
    except AppError as e:
        # Handle our custom application errors (ValidationError, TemplateError)
        logging.warning(f"Application error: {e}")
        return _create_error_response(str(e), e.status_code)
    except Exception as e:
        logging.exception("Unexpected error")
        return _create_error_response(f"Internal server error: {e}", 500)
