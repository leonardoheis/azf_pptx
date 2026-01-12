import logging

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE

from helpers.exceptions import TemplateError
from helpers.utils import (
    _add_bullet,
    _add_bullet_runs,
    _add_section_header,
    _extract_urls,
    _find_shape_with_token,
    _is_url,
    _norm,
    _parse_date,
)

# Font sizes for hierarchical layout
TITLE_FONT_SIZE = 12  # Section headers (e.g., "Latest Relevant News:")
BULLET_FONT_SIZE = 11  # First level bullets (main content)
SUB_BULLET_FONT_SIZE = 10  # Second level bullets (details like "Date published:", "Source:")


# --------------------------------------------------------------------
# Function 3 (generic): {{CompanyResearch3}} -> hierarchical bullets + links
# --------------------------------------------------------------------
def fill_company_research3(prs: Presentation, payload: dict):
    token = "{{CompanyResearch3}}"

    slide, shape = _find_shape_with_token(prs, token)
    if not shape:
        raise TemplateError(f"Token '{token}' not found in any slide")

    if not payload:
        logging.warning("CompanyResearch3 payload is empty; slide will be left blank.")
        return

    # Preserve placeholder geometry to reuse on cloned slides
    base_left, base_top, base_width, base_height = shape.left, shape.top, shape.width, shape.height

    # Gather sections in the original order
    sections = list[tuple](payload.items())

    # Ensure the first slide textbox is ready
    tf_first = shape.text_frame
    tf_first.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf_first.clear()

    def _add_slide_after(prs_obj: Presentation, ref_slide, layout):
        """Add a slide, place it immediately after ref_slide, and clear existing shapes."""
        new_slide = prs_obj.slides.add_slide(layout)
        sldIdLst = prs_obj.slides._sldIdLst  # reorder to desired position
        new_id = sldIdLst[-1]
        sldIdLst.remove(new_id)
        ref_idx = list(prs_obj.slides).index(ref_slide)
        sldIdLst.insert(ref_idx + 1, new_id)
        # Remove any shapes/placeholders so we paste only the new content we add next
        for shp in list(new_slide.shapes):
            new_slide.shapes._spTree.remove(shp._element)
        return new_slide

    # Helper to get or create target slide/shape for each chunk
    def _get_target_tf(chunk_index: int):
        if chunk_index == 0:
            target_slide = slide
            tf_local = tf_first
        else:
            target_slide = _add_slide_after(prs, slide, slide.slide_layout)
            target_shape = target_slide.shapes.add_textbox(base_left, base_top, base_width, base_height)
            tf_local = target_shape.text_frame
            tf_local.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            tf_local.clear()
            logging.info("CompanyResearch3 added continuation slide #%s", chunk_index + 1)
        return target_slide, tf_local

    # -------- internal helpers without external configuration --------
    def _section_items(section_value):
        """
        Extracts the list of items from any section generically:
        - dict with a list (key containing 'list/items/entries/highlights/data/points') -> that list + meta
        - direct list -> items=list, meta=None
        - dict without lists -> [dict] + dict (treated as a single item)
        - primitive -> [{'Value': primitive}] + None
        """
        if isinstance(section_value, dict):
            # look for a "natural" list
            for k, v in section_value.items():
                if isinstance(v, list) and any(
                    w in _norm(k) for w in ["list", "items", "entries", "highlights", "data", "points"]
                ):
                    return v, section_value
            # if none found, take the first list that appears
            for v in section_value.values():
                if isinstance(v, list):
                    return v, section_value
            # flat dict
            return [section_value], section_value
        elif isinstance(section_value, list):
            return section_value, None
        else:
            return [{"Value": section_value}], None

    def _score_main_kv(k, v):
        """
        Scores a key/value pair to choose the main line of an item:
        - Prefers longer strings
        - Gives a small bonus to "title-like" keys (title/name/headline/summary)
        - Penalizes pure URLs
        """
        if isinstance(v, str):
            if _is_url(v):
                return 0.5  # URLs are not good titles
            base = min(len(v.strip()), 200) / 200.0  # normalize by length
        elif isinstance(v, (int, float)):
            base = 0.4
        elif isinstance(v, dict):
            base = 0.3
        elif isinstance(v, list):
            base = 0.35
        else:
            base = 0.2

        nk = _norm(k)
        if any(w in nk for w in ["title", "name", "headline", "subject", "summary", "objective"]):
            base += 0.3
        return base

    def _choose_main_text(item_dict):
        """
        Chooses the main text of an item without depending on fixed names:
        - Maximum score via _score_main_kv
        - If no useful strings, compact as "k: v; ..."
        """
        best_k, best_v, best_score = None, None, -1.0
        for k, v in item_dict.items():
            sc = _score_main_kv(k, v)
            if sc > best_score:
                best_k, best_v, best_score = k, v, sc

        if isinstance(best_v, str) and best_v.strip():
            return best_k, best_v.strip()
        # if the "best" is not a string, try another decent string
        for k, v in item_dict.items():
            if isinstance(v, str) and v.strip() and not _is_url(v):
                return k, v.strip()
        # last resort: compact the dict
        try:
            return best_k, "; ".join(f"{k}: {v}" for k, v in item_dict.items() if v not in (None, ""))
        except Exception:
            # hard fallback
            return best_k, str(next(iter(item_dict.values()), ""))

    def _key_priority(k, v):
        """
        Generic ordering of subfields:
        0: summary/description
        1: dates (date/as of/fiscal)
        2: "normal" values (text/numbers)
        3: URLs and sources
        """
        nk = _norm(k)
        if any(w in nk for w in ["summary", "description", "details", "overview"]):
            return 0
        if any(w in nk for w in ["date", "as of", "fiscal year", "fy"]):
            return 1
        if isinstance(v, str) and _is_url(v):
            return 3
        if any(w in nk for w in ["url", "link", "source", "reference"]):
            return 3
        return 2

    def _order_subkeys(item_dict, main_key_used):
        keys = [k for k in item_dict.keys() if k != main_key_used]
        # sort by priority (+ stable alphabetical)
        return sorted(keys, key=lambda k: (_key_priority(k, item_dict[k]), _norm(k)))

    def _section_suffix_from_meta(meta_dict):
        """If meta contains FY/As Of/Date, add a readable suffix to the header."""
        if not isinstance(meta_dict, dict):
            return ""
        # choose the first reasonable date
        for k, v in meta_dict.items():
            nk = _norm(k)
            if isinstance(v, str) and any(w in nk for w in ["fiscal year", "as of", "date"]):
                nice = _parse_date(v)
                if "fiscal year" in nk or nk == "fy":
                    return f" (FY {nice})" if nice else ""
                return f" ({nice})" if nice else ""
        return ""

    def _emit_value_as_bullets(label, value, level=1):
        """Generic render of a value as bullets/sub-bullets."""
        if value in (None, ""):
            return

        # pure URL
        if isinstance(value, str) and _is_url(value):
            _add_bullet_runs(
                tf, [{"text": f"{label}: "}, {"text": value, "link": value}], level=level, size=SUB_BULLET_FONT_SIZE
            )
            return

        # list
        if isinstance(value, list):
            if all(isinstance(x, (str, int, float)) for x in value):
                for x in value:
                    _add_bullet(tf, f"{label}: {x}", level=level, size=SUB_BULLET_FONT_SIZE)
            else:
                for x in value:
                    if isinstance(x, dict):
                        mk, mv = _choose_main_text(x)
                        _add_bullet(tf, f"{label}: {mv}", level=level, size=SUB_BULLET_FONT_SIZE)
                        # internal URLs
                        for u in _extract_urls(x):
                            _add_bullet_runs(
                                tf,
                                [{"text": "link: "}, {"text": u, "link": u}],
                                level=level + 1,
                                size=SUB_BULLET_FONT_SIZE,
                            )
                    else:
                        _add_bullet(tf, f"{label}: {x}", level=level, size=SUB_BULLET_FONT_SIZE)
            return

        # dict
        if isinstance(value, dict):
            mk, mv = _choose_main_text(value)
            # main line of the sub-dict
            _add_bullet(tf, f"{label}: {mv}", level=level, size=SUB_BULLET_FONT_SIZE)
            # remaining fields of the sub-dict
            for sk in _order_subkeys(value, mk):
                sv = value.get(sk)
                _emit_value_as_bullets(sk, sv, level=level + 1)
            return

        # string/numeric
        nk = _norm(label)
        vtxt = (
            _parse_date(value)
            if isinstance(value, str) and any(s in nk for s in ["date", "as of", "fiscal year", "fy"])
            else str(value)
        )
        _add_bullet(tf, f"{label}: {vtxt}", level=level, size=SUB_BULLET_FONT_SIZE)

    # -------- generic section traversal (in order of appearance) --------
    # Render two section headers per slide (generic, no hardcoded titles)
    for idx in range(0, len(sections), 2):
        _, tf = _get_target_tf(idx // 2)
        for section_name, section_value in sections[idx : idx + 2]:
            items, meta = _section_items(section_value)
            suffix = _section_suffix_from_meta(meta)
            _add_section_header(tf, f"{section_name}{suffix}:", size=TITLE_FONT_SIZE)

            for it in items:
                if isinstance(it, dict):
                    mk, mv = _choose_main_text(it)
                    _add_bullet(tf, mv, level=0, size=BULLET_FONT_SIZE)

                    # subfields of the item
                    for sk in _order_subkeys(it, mk):
                        sv = it.get(sk)
                        _emit_value_as_bullets(sk, sv, level=1)
                elif isinstance(it, list):
                    # list of primitives in an item
                    for x in it:
                        _add_bullet(tf, str(x), level=0, size=BULLET_FONT_SIZE)
                else:
                    # primitive
                    _add_bullet(tf, str(it), level=0, size=BULLET_FONT_SIZE)
